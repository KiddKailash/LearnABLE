"""
Utility module for generating educational content in DOCX, PDF, PPTX, and audio formats.
Provides functions to format plain text into styled documents and convert it to multimedia resources.
"""

import os
import re
import textwrap
import requests
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.platypus import (
    BaseDocTemplate, Frame, PageTemplate, Paragraph,
    Spacer, ListFlowable, ListItem, HRFlowable, Image as RLImage
)
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib import colors


OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")


def create_docx_from_text(text, path, title=None, images=None):
    """
    Create a .docx file from plain text, inferring a title if not provided, stripping unwanted tips/reminders,
    and styling headings ending with a colon.
    """
    # Strip trailing tips/reminders
    body = re.split(r"Tips for Using Tools:", text)[0]
    body = re.split(r"Remember,", body)[0]

    # Normalize to lines & remove [Heading X] prefixes
    lines = [re.sub(r"^\[Heading \d+\]\s*", "", l) for l in body.splitlines()]

    # Infer title from first non-empty line if not given
    if title is None:
        for i, line in enumerate(lines):
            if line.strip():
                title = line.strip()
                lines.pop(i)
                break

    # Build document
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)

    # Add title
    if title:
        p = doc.add_paragraph(title, style='Title')
        p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        doc.add_paragraph()

    # Process lines
    for line in lines:
        stripped = line.strip()
        if not stripped:
            doc.add_paragraph()
            continue

        # Generic heading: ends with ':'
        if stripped.endswith(':') and not stripped[0] in ('-', '*', '•'):
            doc.add_paragraph(stripped.rstrip(':'), style='Heading 2')
            continue

        # Headings: **Heading** or **Heading**:
        m = re.match(r"^\*\*(.+?)\*\*(?::)?$", stripped)
        if m:
            doc.add_paragraph(m.group(1).strip(), style='Heading 2')
            continue

        # Bold label: **Label:** content
        m = re.match(r"^\*\*(.+?)\*\*:\s*(.*)$", stripped)
        if m:
            p = doc.add_paragraph()
            run = p.add_run(f"{m.group(1).strip()}: ")
            run.bold = True
            run.font.size = Pt(12)
            p.add_run(m.group(2).strip())
            continue

        # Bullets: -, *, or •
        m = re.match(r"^[\-\*•]\s+(.*)$", stripped)
        if m:
            doc.add_paragraph(m.group(1).strip(), style='List Bullet')
            continue

        # Regular paragraph
        img_match = re.match(r"\[IMAGE_(\d+)\]", stripped)
        if img_match and images:
            idx = int(img_match.group(1))
            if idx < len(images) and os.path.exists(images[idx]['path']):
                doc.add_picture(images[idx]['path'], width=Inches(5.5))
            continue

        doc.add_paragraph(stripped)

    if images:
        doc.add_page_break()
        doc.add_paragraph("Visual References", style='Heading 2')
        for img in images:
            if os.path.exists(img['path']):
                doc.add_picture(img['path'], width=Inches(5.5))
    doc.save(path)


def create_pdf_from_text(text, path, title=None, images=None):
    """
    Generate a cleanly styled PDF with header line, title, sections, bullets & wrapped text.
    """
    # Strip off any trailing tips/reminders
    body = re.split(r"Tips for Using Tools:", text)[0]
    body = re.split(r"Remember,", body)[0]
    lines = [l for l in body.splitlines() if l.strip()]

    # Infer title if not provided
    if title is None and lines:
        title = lines.pop(0).strip()

    # Build a simple doc template with a single frame
    doc = BaseDocTemplate(path, pagesize=A4,
                          leftMargin=inch*0.75, rightMargin=inch*0.75,
                          topMargin=inch*0.75, bottomMargin=inch*0.75)
    frame = Frame(doc.leftMargin, doc.bottomMargin,
                  doc.width, doc.height, id='normal')
    doc.addPageTemplates([PageTemplate(id='OneCol', frames=frame)])

    # Styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('Title', parent=styles['Title'],
                                 fontName='Helvetica-Bold', fontSize=20,
                                 textColor=colors.HexColor('#2E74B5'),
                                 spaceAfter=12, alignment=1)
    section_style = ParagraphStyle('Section', parent=styles['Heading2'],
                                   fontName='Helvetica-Bold', fontSize=14,
                                   textColor=colors.HexColor('#2E74B5'),
                                   spaceBefore=12, spaceAfter=6)
    body_style = ParagraphStyle('Body', parent=styles['BodyText'],
                                fontName='Helvetica', fontSize=11,
                                leading=14)

    elements = []
    # Draw a top rule line
    elements.append(HRFlowable(width='100%', thickness=1, color=colors.HexColor(
        '#2E74B5'), spaceBefore=0, spaceAfter=12))
    # Title
    if title:
        elements.append(Paragraph(title, title_style))

    # Process remaining lines
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        # Section headings marked **Heading**
        m = re.match(r"^\*\*(.+?)\*\*(?::)?$", line)
        if m:
            elements.append(Paragraph(m.group(1).strip(), section_style))
            i += 1
            continue

        # Bold label paragraphs **Label:** content
        m = re.match(r"^\*\*(.+?)\*\*:\s*(.*)$", line)
        if m:
            label, content = m.groups()
            wrapped = textwrap.wrap(f"<b>{label}:</b> {content}", width=100)
            for w in wrapped:
                elements.append(Paragraph(w, body_style))
            i += 1
            continue

        # Bullet lists if line starts with -, *, or •
        if line and line[0] in ('-', '*', '•'):
            items = []
            while i < len(lines) and lines[i].strip()[0] in ('-', '*', '•'):
                item_text = lines[i].strip()[1:].strip()
                items.append(ListItem(Paragraph(item_text, body_style)))
                i += 1
            elements.append(ListFlowable(items, bulletType='bullet',
                            leftIndent=12, spaceBefore=4, spaceAfter=4))
            continue

        # Regular paragraph
        wrapped = textwrap.wrap(line, width=100)
        for w in wrapped:
            elements.append(Paragraph(w, body_style))
        i += 1

    if images:
        elements.append(Spacer(1, 12))
        elements.append(Paragraph("Visual References", section_style))
        for img in images:
            if os.path.exists(img['path']):
                rl_img = RLImage(img['path'], width=4*inch, height=3*inch)
                rl_img.hAlign = 'CENTER'
                elements.append(Spacer(1, 12))
                elements.append(rl_img)

    doc.build(elements)


def create_pptx_from_text(slide_pairs, path):
    """
    slide_pairs: list of tuples (title: str, content: str, images: list).
                 images may be dicts with keys 'path','left','top','width','height'
                 or plain string filepaths.
    path:       output .pptx filepath.
    """
    prs = Presentation()
    title_font_size = Pt(36)
    content_font_size = Pt(20)
    all_images = []

    # 1) Build content slides and collect all images
    for idx, (title, content, images) in enumerate(slide_pairs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # — Title box —
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(0.3), Inches(9), Inches(1)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = RGBColor(91, 155, 213)
        tb = title_box.text_frame
        tb.text = title.strip()
        p = tb.paragraphs[0]
        p.font.size = title_font_size
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # — Content box —
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.5), Inches(8), Inches(4)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(content.split("\n")):
            if not line.strip():
                continue
            paragraph = tf.add_paragraph() if i else tf.paragraphs[0]
            paragraph.text = line.strip()
            paragraph.level = 0
            paragraph.font.size = content_font_size
            paragraph.font.name = "Calibri"
            paragraph.font.color.rgb = RGBColor(50, 50, 50)

        # Collect any images for the "Visual References" slides
        if images:
            all_images.extend(images)

    # 2) After all content slides are done, generate the Visual References slides
    if all_images:
        def add_image_slide():
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(0.3), Inches(9), Inches(1)
            )
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = RGBColor(91, 155, 213)
            tf = title_box.text_frame
            tf.text = "Visual References"
            p = tf.paragraphs[0]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            return slide

        # layout parameters
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        image_width = Inches(6.5)
        image_height = Inches(4.5)
        top_margin = Inches(1.4)
        spacing = Inches(0.6)

        # start first image slide
        img_slide = add_image_slide()
        y_offset = top_margin

        for img in all_images:
            img_path = img.get("path") if isinstance(img, dict) else img
            if not (img_path and os.path.exists(img_path)):
                continue

            # only create a new slide if the next image won't fit
            if y_offset + image_height > slide_height - Inches(0.5):
                img_slide = add_image_slide()
                y_offset = top_margin

            # center the image horizontally
            left = (slide_width - image_width) // 2
            img_slide.shapes.add_picture(
                img_path,
                left=left,
                top=y_offset,
                width=image_width,
                height=image_height
            )
            y_offset += image_height + spacing

    # 3) Save the presentation
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)


# nova – female, natural and calm

# shimmer – female, bright and expressive

# echo – male, neutral and smooth

# onyx – male, deeper tone

# fable – male, slightly theatrical

def create_audio_from_text(text, path, voice="nova", speed=0.95):
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True)

        url = "https://api.openai.com/v1/audio/speech"
        headers = {
            "Authorization": f"Bearer {OPENAI_API_KEY}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "tts-1",
            "voice": voice,
            "input": text,
            "speed": speed
        }

        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()

        with open(path, "wb") as f:
            f.write(response.content)

        print(f"[AUDIO] Audio saved at {path}")
        return True
    except Exception as e:
        print(f"[AUDIO ERROR] {e}")
        return False
