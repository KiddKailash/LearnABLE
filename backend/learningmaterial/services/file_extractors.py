"""
Utility functions to extract text and image data from PDF, DOCX, and PPTX files.
These help standardize content input for further AI processing or transformation.
"""

import fitz  # PyMuPDF to read pdfs and extract text
from docx import Document
from pptx import Presentation
import os
import tempfile
from pptx.enum.shapes import PP_PLACEHOLDER
from docx.opc.constants import RELATIONSHIP_TYPE as RT

tmp_dir = os.path.join(tempfile.gettempdir(), "pptx_images")
os.makedirs(tmp_dir, exist_ok=True)

def extract_text_from_pdf(path):
    """
    Extract text and images from a PDF file.
    Returns a tuple: (text, images)
    """
    lines = []
    images = []
    with fitz.open(path) as doc:
        for page_index in range(len(doc)):
            page = doc[page_index]
            # Get text
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        sentence = " ".join([span["text"]
                                            for span in line["spans"]])
                        lines.append(sentence)

            # Get images
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_path = os.path.join(
                    tmp_dir, f"pdf_img_{page_index}_{img_index}.{image_ext}")
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                images.append({'path': image_path})

    return "\n".join(lines), images


def extract_text_from_docx(path):
    """
    Extract text and image paths from a DOCX file.
    Returns a tuple: (text, images)
    """
    doc = Document(path)
    lines = []
    images = []

    # Text extraction
    for para in doc.paragraphs:
        style = para.style.name.lower()
        text = para.text.strip()
        if not text:
            continue
        if "heading" in style:
            lines.append(f"[{para.style.name}] {text}")
        else:
            lines.append(text)

    # Image extraction
    rels = doc.part._rels
    for rel in rels.values():
        if rel.reltype == RT.IMAGE:
            image_part = rel.target_part
            image_data = image_part.blob
            image_ext = image_part.content_type.split(
                "/")[-1]  # e.g., image/png
            img_path = os.path.join(
                tmp_dir, f"docx_img_{len(images)}.{image_ext}")
            with open(img_path, "wb") as f:
                f.write(image_data)
            images.append({'path': img_path})

    return "\n".join(lines), images


def extract_text_from_pptx(path):
    """
    Extract text and image metadata from a PowerPoint file.
    Returns a list of slides with title, content, and image data.
    """
    prs = Presentation(path)
    slide_data = []

    for idx, slide in enumerate(prs.slides):
        title = ""
        content = ""
        images = []

        for shape in slide.shapes:
            # 1) Capture ALL text frames, not just placeholders
            if hasattr(shape, "text_frame") and shape.text_frame:
                # Join all paragraphs in this shape
                text = "\n".join(p.text.strip()
                                 for p in shape.text_frame.paragraphs if p.text.strip())
                if not text:
                    pass
                elif shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    title = text
                else:
                    # Anything else (body placeholders or regular text boxes) goes into content
                    content += text + "\n"

            # 2) Capture images as before
            if hasattr(shape, "image"):
                blob = shape.image.blob
                ext = shape.image.ext
                fname = f"slide_{idx}_img_{len(images)}.{ext}"
                fpath = os.path.join(tmp_dir, fname)
                with open(fpath, "wb") as f:
                    f.write(blob)
                images.append({
                    "path":   fpath,
                    "left":   shape.left,
                    "top":    shape.top,
                    "width":  shape.width,
                    "height": shape.height
                })

        slide_data.append({
            "title":   title,
            "content": content.strip(),
            "images":  images
        })

    return slide_data
